"""
Post-process the generated .pptx to add:

  1. Subtle "fade" slide transitions on every slide.

  2. Auto-playing per-slide entrance animations:
       - Visual blocks (a card with its inner ellipse/icon/text) are
         WRAPPED into a <p:grpSp> group shape, then animated as a SINGLE
         "Appear" effect. PowerPoint treats the whole group as one unit,
         which is what makes each row enter at once rather than its
         sub-elements cascading one after another.
       - A large right-side photo (when present) "flies in" from the
         right AFTER the body blocks.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

PATH = "Internet-i-Intranet.pptx"

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "{%s}" % P_NS
A = "{%s}" % A_NS

EMU = 914400


def _q(tag):
    return P + tag


# ============================================================
# Slide transition
# ============================================================
def add_transition(slide):
    root = slide.element
    existing = root.find(_q("transition"))
    if existing is not None:
        root.remove(existing)
    trans = etree.Element(_q("transition"))
    trans.set("spd", "med")
    etree.SubElement(trans, _q("fade"))
    cSld = root.find(_q("cSld"))
    idx = list(root).index(cSld) + 1
    nxt = root[idx] if idx < len(root) else None
    if nxt is not None and nxt.tag == _q("clrMapOvr"):
        idx += 1
    root.insert(idx, trans)


# ============================================================
# Shape filtering
# ============================================================
def is_animatable(shape):
    """Skip decorations: backgrounds, headers, footers, accent strips,
    decorative blobs. Animate everything else in the body area."""
    if shape.top is None or shape.left is None:
        return False
    if shape.width is None or shape.height is None:
        return False
    if shape.top < 0.5 * EMU or shape.top > 5.0 * EMU:
        return False
    if shape.width > 9.0 * EMU:
        return False
    if shape.width > 3.4 * EMU and shape.height > 3.4 * EMU:
        return False
    if shape.height < 0.06 * EMU:
        return False
    return True


def is_hero_image(shape):
    """The prominent right-side photo on content slides."""
    if shape.shape_type != 13:  # PICTURE
        return False
    if shape.width is None or shape.height is None:
        return False
    if shape.width < 2.5 * EMU or shape.height < 1.5 * EMU:
        return False
    return True


def y_center(shape):
    return (shape.top or 0) + (shape.height or 0) // 2


def get_animation_clusters(slide):
    """Group body shapes into visual blocks so each animation step
    corresponds to one logical card / row / column.

    Strategy:
      1. Detect *large containers* — auto-shapes ≥3"×1.8". These are
         the "tall cards" (slide 11 columns, slide 12 conclusion cards)
         whose contents span a wide Y range and would otherwise be
         split by row-clustering. All shapes whose center falls inside
         a container go in that container's group.
      2. Remaining shapes are clustered by vertical-center Y position
         (rows), then any wide Y-cluster containing multiple side-by-
         side card backgrounds is split into one sub-cluster per card
         (handles 2x2 grids on slide 4 / slide 8).
      3. Container clusters and row clusters are merged and sorted
         top-to-bottom.
      4. Hero photo (if present) is appended as the final group, with
         a Fly-In-from-right preset.
    """
    body = []
    hero = None
    for shape in slide.shapes:
        if not is_animatable(shape):
            continue
        if is_hero_image(shape):
            if hero is None or (shape.width * shape.height) > (hero.width * hero.height):
                if hero is not None:
                    body.append(hero)
                hero = shape
            else:
                body.append(shape)
        else:
            body.append(shape)

    # --- Step 1: detect large containers ---
    # A "container" is a sizeable empty rectangle/roundRect (no text).
    # Excluding text-bearing shapes prevents e.g. a tall bullet-list
    # text box from being mistaken for a card background.
    containers = []
    for s in body:
        if s.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if not s.width or not s.height:
            continue
        if s.width <= 3.0 * EMU or s.height <= 1.8 * EMU:
            continue
        # Skip if the shape has visible text — it's a content shape, not a container
        if s.has_text_frame and s.text_frame.text.strip():
            continue
        containers.append(s)

    # --- Step 2: assign shapes to their innermost container ---
    container_groups = {id(c): [c] for c in containers}
    not_contained = []
    for s in body:
        if s in containers:
            continue
        sx = s.left + (s.width or 0) // 2
        sy = s.top + (s.height or 0) // 2
        # Pick the smallest container that contains this shape's center
        assigned = None
        for c in containers:
            if (c.left <= sx <= c.left + c.width
                    and c.top <= sy <= c.top + c.height):
                if assigned is None or (c.width * c.height) < (assigned.width * assigned.height):
                    assigned = c
        if assigned is not None:
            container_groups[id(assigned)].append(s)
        else:
            not_contained.append(s)

    # --- Step 3: Y-cluster + X-split the remaining shapes ---
    not_contained.sort(key=y_center)
    row_clusters = []
    current = []
    current_yc = None
    Y_TOL = 0.5 * EMU
    for s in not_contained:
        yc = y_center(s)
        if current_yc is None or (yc - current_yc) < Y_TOL:
            current.append(s)
            if current_yc is None:
                current_yc = yc
        else:
            row_clusters.append(current)
            current = [s]
            current_yc = yc
    if current:
        row_clusters.append(current)

    expanded_rows = []
    for cl in row_clusters:
        for sub in split_cluster_by_cards(cl):
            if sub:
                expanded_rows.append(sub)

    # --- Step 4: combine container clusters + row clusters, sort by Y ---
    all_clusters = []
    for c in containers:
        shapes_in = container_groups[id(c)]
        min_y = min((s.top or 0) + (s.height or 0) // 2 for s in shapes_in)
        all_clusters.append((min_y, shapes_in))
    for sub in expanded_rows:
        min_y = min((s.top or 0) + (s.height or 0) // 2 for s in sub)
        all_clusters.append((min_y, sub))
    all_clusters.sort(key=lambda x: x[0])

    result = [{"preset_id": 1, "preset_subtype": 0, "shapes": cl}
              for _, cl in all_clusters]
    if hero is not None:
        # Photo uses the same simple "Appear" preset for max compatibility
        # across PowerPoint versions, web, mobile, Keynote, LibreOffice.
        # (Fly-In with ppt_x formulas is PowerPoint-specific and may
        # render inconsistently or be ignored on other renderers.)
        result.append({"preset_id": 1, "preset_subtype": 0, "shapes": [hero]})
    return result


def split_cluster_by_cards(cluster):
    """If a Y-cluster contains 2+ card-sized backgrounds side by side
    (e.g. slide 4's 2x2 grid), split it horizontally so each card +
    its inner shapes animates independently."""
    cards = []
    for s in cluster:
        if s.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if not s.width or not s.height:
            continue
        # A "card background" is a sizeable auto-shape — large enough that
        # multiple of them in the same row imply distinct cards.
        if s.width > 1.5 * EMU and s.height > 0.5 * EMU:
            cards.append(s)

    if len(cards) <= 1:
        return [cluster]

    cards.sort(key=lambda s: s.left)
    sub_clusters = [[] for _ in cards]
    for s in cluster:
        sx_center = s.left + (s.width or 0) // 2
        # Assign to the card whose X-range contains this shape's center
        assigned = False
        for i, c in enumerate(cards):
            if c.left <= sx_center <= c.left + c.width:
                sub_clusters[i].append(s)
                assigned = True
                break
        if not assigned:
            distances = [abs(sx_center - (c.left + c.width // 2)) for c in cards]
            min_idx = distances.index(min(distances))
            sub_clusters[min_idx].append(s)
    return sub_clusters


# ============================================================
# <p:grpSp> wrapping — bind a cluster's shapes into a single group
# so PowerPoint animates them as one entity.
# ============================================================
def next_shape_id(slide_root):
    max_id = 0
    for cNvPr in slide_root.iter(P + "cNvPr"):
        try:
            max_id = max(max_id, int(cNvPr.get("id")))
        except (TypeError, ValueError):
            pass
    return max_id + 1


def wrap_in_group(slide, shapes, name):
    """Wrap shape XML elements inside a new <p:grpSp>. Returns the
    group's shape ID. The group's bounding box is set to enclose all
    children, with chOff/chExt matching off/ext so children render
    at their original absolute positions."""
    slide_root = slide.element
    spTree = slide_root.find(P + "cSld").find(P + "spTree")

    lefts   = [s.left for s in shapes]
    tops    = [s.top for s in shapes]
    rights  = [s.left + s.width for s in shapes]
    bottoms = [s.top + s.height for s in shapes]
    g_x = min(lefts)
    g_y = min(tops)
    g_w = max(rights) - g_x
    g_h = max(bottoms) - g_y

    new_id = next_shape_id(slide_root)
    shape_els = [s._element for s in shapes]
    children = list(spTree)
    try:
        # Sort by current spTree position to preserve z-order inside the group
        shape_els.sort(key=lambda el: children.index(el))
        insert_at = children.index(shape_els[0])
    except ValueError:
        return None

    grpSp = etree.Element(P + "grpSp")
    nvGrpSpPr = etree.SubElement(grpSp, P + "nvGrpSpPr")
    etree.SubElement(nvGrpSpPr, P + "cNvPr", id=str(new_id), name=name)
    etree.SubElement(nvGrpSpPr, P + "cNvGrpSpPr")
    etree.SubElement(nvGrpSpPr, P + "nvPr")

    grpSpPr = etree.SubElement(grpSp, P + "grpSpPr")
    xfrm = etree.SubElement(grpSpPr, A + "xfrm")
    etree.SubElement(xfrm, A + "off",   x=str(g_x), y=str(g_y))
    etree.SubElement(xfrm, A + "ext",   cx=str(g_w), cy=str(g_h))
    etree.SubElement(xfrm, A + "chOff", x=str(g_x), y=str(g_y))
    etree.SubElement(xfrm, A + "chExt", cx=str(g_w), cy=str(g_h))

    for el in shape_els:
        el.getparent().remove(el)
        grpSp.append(el)

    spTree.insert(insert_at, grpSp)
    return new_id


# ============================================================
# Timing XML
# ============================================================
class IdGen:
    def __init__(self, start=3):
        self.n = start

    def next(self):
        v = self.n
        self.n += 1
        return v


def _set_visibility(parent, sid, gen):
    st = etree.SubElement(parent, _q("set"))
    cBhvr = etree.SubElement(st, _q("cBhvr"))
    cTn = etree.SubElement(cBhvr, _q("cTn"), id=str(gen.next()), dur="1", fill="hold")
    stc = etree.SubElement(cTn, _q("stCondLst"))
    etree.SubElement(stc, _q("cond"), delay="0")
    tgtEl = etree.SubElement(cBhvr, _q("tgtEl"))
    etree.SubElement(tgtEl, _q("spTgt"), spid=str(sid))
    attrLst = etree.SubElement(cBhvr, _q("attrNameLst"))
    attrName = etree.SubElement(attrLst, _q("attrName"))
    attrName.text = "style.visibility"
    to = etree.SubElement(st, _q("to"))
    etree.SubElement(to, _q("strVal"), val="visible")


def _flyin_from_right(parent, sid, gen, duration_ms=500):
    anim = etree.SubElement(parent, _q("anim"), calcmode="lin", valueType="num")
    cBhvr = etree.SubElement(anim, _q("cBhvr"), additive="base")
    etree.SubElement(cBhvr, _q("cTn"),
                     id=str(gen.next()), dur=str(duration_ms), fill="hold")
    tgtEl = etree.SubElement(cBhvr, _q("tgtEl"))
    etree.SubElement(tgtEl, _q("spTgt"), spid=str(sid))
    attrLst = etree.SubElement(cBhvr, _q("attrNameLst"))
    attrName = etree.SubElement(attrLst, _q("attrName"))
    attrName.text = "ppt_x"
    tavLst = etree.SubElement(anim, _q("tavLst"))
    tav1 = etree.SubElement(tavLst, _q("tav"), tm="0")
    etree.SubElement(etree.SubElement(tav1, _q("val")), _q("strVal"),
                     val="1+#ppt_w/2")
    tav2 = etree.SubElement(tavLst, _q("tav"), tm="100000")
    etree.SubElement(etree.SubElement(tav2, _q("val")), _q("strVal"),
                     val="#ppt_x")


def build_timing(specs, stagger_ms=350, hero_extra_ms=200):
    """Build a <p:timing> element. `specs` is a list of dicts:
       { shape_id: int, preset_id: int, preset_subtype: int }"""
    gen = IdGen(start=3)

    timing = etree.Element(_q("timing"))
    tnLst = etree.SubElement(timing, _q("tnLst"))
    par_root = etree.SubElement(tnLst, _q("par"))
    cTn_root = etree.SubElement(par_root, _q("cTn"),
                                id="1", dur="indefinite",
                                restart="never", nodeType="tmRoot")
    children_root = etree.SubElement(cTn_root, _q("childTnLst"))
    seq = etree.SubElement(children_root, _q("seq"),
                           concurrent="1", nextAc="seek")
    cTn_seq = etree.SubElement(seq, _q("cTn"),
                               id="2", dur="indefinite", nodeType="mainSeq")
    children_seq = etree.SubElement(cTn_seq, _q("childTnLst"))

    par_step = etree.SubElement(children_seq, _q("par"))
    cTn_step = etree.SubElement(par_step, _q("cTn"),
                                id=str(gen.next()), fill="hold")
    stc_step = etree.SubElement(cTn_step, _q("stCondLst"))
    etree.SubElement(stc_step, _q("cond"), delay="0")  # auto-start
    children_step = etree.SubElement(cTn_step, _q("childTnLst"))

    par_inner = etree.SubElement(children_step, _q("par"))
    cTn_inner = etree.SubElement(par_inner, _q("cTn"),
                                 id=str(gen.next()), fill="hold")
    stc_inner = etree.SubElement(cTn_inner, _q("stCondLst"))
    etree.SubElement(stc_inner, _q("cond"), delay="0")
    children_inner = etree.SubElement(cTn_inner, _q("childTnLst"))

    animated_ids = []
    is_first = True
    for spec in specs:
        sid = spec["shape_id"]
        preset_id = spec["preset_id"]
        preset_subtype = spec["preset_subtype"]

        # Delay (relative to previous effect's start, for withEffect)
        if is_first:
            delay = 0
        else:
            delay = stagger_ms

        animated_ids.append(sid)

        par_eff = etree.SubElement(children_inner, _q("par"))
        cTn_eff = etree.SubElement(par_eff, _q("cTn"),
            id=str(gen.next()),
            presetID=str(preset_id), presetClass="entr",
            presetSubtype=str(preset_subtype),
            fill="hold", grpId="0",
            nodeType="clickEffect" if is_first else "withEffect",
        )
        stc_eff = etree.SubElement(cTn_eff, _q("stCondLst"))
        etree.SubElement(stc_eff, _q("cond"), delay=str(delay))
        children_eff = etree.SubElement(cTn_eff, _q("childTnLst"))

        _set_visibility(children_eff, sid, gen)
        if preset_id == 2:  # Fly In from right
            _flyin_from_right(children_eff, sid, gen, duration_ms=500)

        is_first = False

    prevCondLst = etree.SubElement(seq, _q("prevCondLst"))
    cond_prev = etree.SubElement(prevCondLst, _q("cond"),
                                 evt="onPrev", delay="0")
    etree.SubElement(etree.SubElement(cond_prev, _q("tgtEl")), _q("sldTgt"))
    nextCondLst = etree.SubElement(seq, _q("nextCondLst"))
    cond_next = etree.SubElement(nextCondLst, _q("cond"),
                                 evt="onNext", delay="0")
    etree.SubElement(etree.SubElement(cond_next, _q("tgtEl")), _q("sldTgt"))

    bldLst = etree.SubElement(timing, _q("bldLst"))
    for sid in animated_ids:
        etree.SubElement(bldLst, _q("bldP"),
                         spid=str(sid), grpId="0", animBg="1")

    return timing


def add_timing_el(slide, timing_el):
    root = slide.element
    existing = root.find(_q("timing"))
    if existing is not None:
        root.remove(existing)
    extLst = root.find(_q("extLst"))
    if extLst is not None:
        root.insert(list(root).index(extLst), timing_el)
    else:
        root.append(timing_el)


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation(PATH)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        add_transition(slide)
        clusters = get_animation_clusters(slide)
        if not clusters:
            continue

        # For each cluster: if >1 shape, wrap in a group and animate the
        # group (a single Appear effect for the whole block).
        specs = []
        for ci, c in enumerate(clusters):
            shapes = c["shapes"]
            if len(shapes) > 1:
                gid = wrap_in_group(slide, shapes,
                                    name=f"Block_s{slide_idx}_c{ci}")
                if gid is None:
                    gid = shapes[0].shape_id
                specs.append({
                    "shape_id": gid,
                    "preset_id": c["preset_id"],
                    "preset_subtype": c["preset_subtype"],
                })
            else:
                specs.append({
                    "shape_id": shapes[0].shape_id,
                    "preset_id": c["preset_id"],
                    "preset_subtype": c["preset_subtype"],
                })

        timing = build_timing(specs)
        add_timing_el(slide, timing)

    prs.save(PATH)
    print(f"✓ Tranzicije + grupisane animacije dodate: {PATH}")


if __name__ == "__main__":
    main()
